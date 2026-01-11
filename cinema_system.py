import json
import os
from datetime import datetime, timedelta
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
import openpyxl
from openpyxl.cell import Cell, MergedCell
from openpyxl.chart import BarChart, Reference
from openpyxl.styles import Font, Alignment, PatternFill
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptxRGBColor


class CinemaSystem:
    def __init__(self):
        self.theatres_dir = "theatres"
        self.reports_dir = "reports"
        if not os.path.exists(self.theatres_dir):
            os.makedirs(self.theatres_dir)
        if not os.path.exists(self.reports_dir):
            os.makedirs(self.reports_dir)

    def add_theatre(self, name):
        theatre_data = {
            "name": name,
            "halls": []
        }
        filename = os.path.join(self.theatres_dir, f"{name}.json")

        if os.path.exists(filename):
            print(f"Кинотеатр '{name}' уже существует!")
            return False

        with open(filename, 'w', encoding='utf-8') as f:
            json.dump(theatre_data, f, ensure_ascii=False, indent=2)
        print(f"Кинотеатр '{name}' успешно добавлен!")
        return True

    def get_theatre(self, name):
        filename = os.path.join(self.theatres_dir, f"{name}.json")
        if not os.path.exists(filename):
            return None

        with open(filename, 'r', encoding='utf-8') as f:
            return json.load(f)

    def save_theatre(self, name, data):
        filename = os.path.join(self.theatres_dir, f"{name}.json")
        with open(filename, 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=2)

    def list_theatres(self):
        files = os.listdir(self.theatres_dir)
        theatres = [f.replace('.json', '') for f in files if f.endswith('.json')]
        return theatres

    def add_hall(self, theatre_name, hall_number, rows, seats_per_row):
        theatre = self.get_theatre(theatre_name)
        if not theatre:
            print(f"Кинотеатр '{theatre_name}' не найден!")
            return False

        for hall in theatre["halls"]:
            if hall["number"] == hall_number:
                print(f"Зал №{hall_number} уже существует в кинотеатре '{theatre_name}'!")
                return False

        hall_data = {
            "number": hall_number,
            "rows": rows,
            "seats_per_row": seats_per_row,
            "sessions": []
        }

        theatre["halls"].append(hall_data)
        self.save_theatre(theatre_name, theatre)
        print(f"Зал №{hall_number} добавлен в кинотеатр '{theatre_name}'!")
        return True

    def create_session(self, theatre_name, hall_number, movie_name, start_time, duration):
        theatre = self.get_theatre(theatre_name)
        if not theatre:
            print(f"Кинотеатр '{theatre_name}' не найден!")
            return False

        hall = None
        for h in theatre["halls"]:
            if h["number"] == hall_number:
                hall = h
                break

        if not hall:
            print(f"Зал №{hall_number} не найден в кинотеатре '{theatre_name}'!")
            return False

        seats = [[False for _ in range(hall["seats_per_row"])] for _ in range(hall["rows"])]

        session_data = {
            "movie": movie_name,
            "start_time": start_time,
            "duration": duration,
            "seats": seats
        }

        hall["sessions"].append(session_data)
        self.save_theatre(theatre_name, theatre)
        print(f"Сеанс фильма '{movie_name}' создан на {start_time}!")
        return True

    def sell_ticket(self, theatre_name, hall_number, session_index, row, seat):
        theatre = self.get_theatre(theatre_name)
        if not theatre:
            print(f"Кинотеатр '{theatre_name}' не найден!")
            return False

        hall = None
        for h in theatre["halls"]:
            if h["number"] == hall_number:
                hall = h
                break

        if not hall:
            print(f"Зал №{hall_number} не найден!")
            return False

        if session_index >= len(hall["sessions"]):
            print("Сеанс не найден!")
            return False

        session = hall["sessions"][session_index]

        if row < 0 or row >= hall["rows"]:
            print(f"Ряд {row + 1} не существует!")
            return False

        if seat < 0 or seat >= hall["seats_per_row"]:
            print(f"Место {seat + 1} не существует!")
            return False

        if session["seats"][row][seat]:
            print(f"Место {row + 1}-{seat + 1} уже занято!")
            return False

        session["seats"][row][seat] = True
        self.save_theatre(theatre_name, theatre)
        print(f"Билет продан! Кинотеатр: {theatre_name}, Зал: {hall_number}, "
              f"Фильм: {session['movie']}, Время: {session['start_time']}, "
              f"Место: Ряд {row + 1}, Место {seat + 1}")
        return True

    def find_nearest_session(self, movie_name):
        current_time = datetime.now()
        nearest_session = None
        nearest_time = None
        nearest_info = None

        for theatre_name in self.list_theatres():
            theatre = self.get_theatre(theatre_name)

            for hall in theatre["halls"]:
                for session_index, session in enumerate(hall["sessions"]):
                    if session["movie"] == movie_name:
                        has_free_seats = False
                        for row in session["seats"]:
                            if False in row:
                                has_free_seats = True
                                break

                        if has_free_seats:
                            try:
                                session_time = datetime.strptime(session["start_time"], "%Y-%m-%d %H:%M")

                                if session_time > current_time:
                                    if nearest_time is None or session_time < nearest_time:
                                        nearest_time = session_time
                                        nearest_session = session
                                        nearest_info = {
                                            "theatre": theatre_name,
                                            "hall": hall["number"],
                                            "session_index": session_index
                                        }
                            except ValueError:
                                pass

        if nearest_session:
            print(f"\nБлижайший сеанс фильма '{movie_name}':")
            print(f"Кинотеатр: {nearest_info['theatre']}")
            print(f"Зал: {nearest_info['hall']}")
            print(f"Время: {nearest_session['start_time']}")
            print(f"Длительность: {nearest_session['duration']} мин")
            return nearest_info
        else:
            print(f"Сеансы фильма '{movie_name}' со свободными местами не найдены.")
            return None

    def print_hall_plan(self, theatre_name, hall_number, session_index):
        theatre = self.get_theatre(theatre_name)
        if not theatre:
            print(f"Кинотеатр '{theatre_name}' не найден!")
            return False

        hall = None
        for h in theatre["halls"]:
            if h["number"] == hall_number:
                hall = h
                break

        if not hall:
            print(f"Зал №{hall_number} не найден!")
            return False

        if session_index >= len(hall["sessions"]):
            print("Сеанс не найден!")
            return False

        session = hall["sessions"][session_index]

        print(f"\n{'=' * 60}")
        print(f"ПЛАН ЗАЛА - Кинотеатр: {theatre_name}, Зал: {hall_number}")
        print(f"Фильм: {session['movie']}, Время: {session['start_time']}")
        print(f"{'=' * 60}")
        print("\n                    ЭКРАН")
        print(f"{'-' * 60}\n")

        print("     ", end="")
        for seat_num in range(hall["seats_per_row"]):
            print(f"{seat_num + 1:3}", end=" ")
        print("\n")

        free_count = 0
        occupied_count = 0

        for row_num, row in enumerate(session["seats"]):
            print(f"Ряд {row_num + 1:2} ", end="")
            for seat in row:
                if seat:
                    print(" X ", end=" ")
                    occupied_count += 1
                else:
                    print(" O ", end=" ")
                    free_count += 1
            print()

        print(f"\n{'=' * 60}")
        print(f"Обозначения: O - свободно, X - занято")
        print(f"Свободных мест: {free_count}, Занятых мест: {occupied_count}")
        print(f"{'=' * 60}\n")
        return True

    def generate_monthly_schedule_docx(self):
        """Генерация расписания сеансов за прошедший месяц в формате DOCX"""
        print("\n--- ГЕНЕРАЦИЯ РАСПИСАНИЯ СЕАНСОВ ---")

        current_date = datetime.now()
        last_month_start = (current_date.replace(day=1) - timedelta(days=1)).replace(day=1)
        last_month_end = current_date.replace(day=1) - timedelta(days=1)

        doc = Document()

        # Заголовок
        title = doc.add_heading('РАСПИСАНИЕ СЕАНСОВ', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        subtitle = doc.add_paragraph(
            f'За период: {last_month_start.strftime("%d.%m.%Y")} - {last_month_end.strftime("%d.%m.%Y")}')
        subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
        subtitle_format = subtitle.runs[0]
        subtitle_format.font.size = Pt(14)
        subtitle_format.font.color.rgb = RGBColor(128, 128, 128)

        doc.add_paragraph()

        # Собираем данные о сеансах
        sessions_found = False

        for theatre_name in self.list_theatres():
            theatre = self.get_theatre(theatre_name)
            theatre_sessions = []

            for hall in theatre["halls"]:
                for session in hall["sessions"]:
                    try:
                        session_time = datetime.strptime(session["start_time"], "%Y-%m-%d %H:%M")
                        if last_month_start <= session_time <= last_month_end:
                            theatre_sessions.append({
                                "hall": hall["number"],
                                "movie": session["movie"],
                                "time": session_time,
                                "duration": session["duration"]
                            })
                    except ValueError:
                        pass

            if theatre_sessions:
                sessions_found = True

                # Заголовок кинотеатра
                theatre_heading = doc.add_heading(f'Кинотеатр: {theatre_name}', 1)
                theatre_heading.runs[0].font.color.rgb = RGBColor(0, 51, 102)

                # Сортируем сеансы по дате
                theatre_sessions.sort(key=lambda x: x["time"])

                # Группируем по датам
                sessions_by_date = {}
                for session in theatre_sessions:
                    date_key = session["time"].strftime("%d.%m.%Y")
                    if date_key not in sessions_by_date:
                        sessions_by_date[date_key] = []
                    sessions_by_date[date_key].append(session)

                # Выводим по датам
                for date_key in sorted(sessions_by_date.keys()):
                    date_heading = doc.add_heading(f'{date_key}', 2)
                    date_heading.runs[0].font.color.rgb = RGBColor(51, 102, 153)

                    for session in sessions_by_date[date_key]:
                        session_text = doc.add_paragraph()
                        session_text.add_run(f'🎬 {session["movie"]}').bold = True
                        session_text.add_run(f'\n   Зал: {session["hall"]} | ')
                        session_text.add_run(f'Время: {session["time"].strftime("%H:%M")} | ')
                        session_text.add_run(f'Длительность: {session["duration"]} мин')
                        session_text.paragraph_format.left_indent = Inches(0.25)

                    doc.add_paragraph()

                doc.add_page_break()

        if not sessions_found:
            doc.add_paragraph("За прошедший месяц сеансы не найдены.")

        # Сохраняем документ
        filename = os.path.join(self.reports_dir, f'schedule_{last_month_start.strftime("%Y-%m")}.docx')
        doc.save(filename)
        print(f"\nРасписание успешно сохранено: {filename}")
        return filename

    def generate_occupancy_chart_xlsx(self):
        """Генерация графика загруженности кинотеатров в зависимости от времени суток в формате XLSX"""
        print("\n--- ГЕНЕРАЦИЯ ГРАФИКА ЗАГРУЖЕННОСТИ ---")

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Загруженность"

        # Заголовок
        ws['A1'] = 'ГРАФИК ЗАГРУЖЕННОСТИ КИНОТЕАТРОВ ПО ВРЕМЕНИ СУТОК'
        ws['A1'].font = Font(size=16, bold=True, color='FFFFFF')
        ws['A1'].fill = PatternFill(start_color='0033AA', end_color='0033AA', fill_type='solid')
        ws['A1'].alignment = Alignment(horizontal='center', vertical='center')
        ws.merge_cells('A1:E1')
        ws.row_dimensions[1].height = 30

        # Временные интервалы
        time_intervals = {
            "Утро (6-12)": (6, 12),
            "День (12-18)": (12, 18),
            "Вечер (18-22)": (18, 22),
            "Ночь (22-6)": (22, 6)
        }

        # Заголовки столбцов
        ws['A3'] = 'Кинотеатр'
        ws['A3'].font = Font(bold=True, size=12)
        ws['A3'].fill = PatternFill(start_color='CCE5FF', end_color='CCE5FF', fill_type='solid')

        col = 2
        for interval_name in time_intervals.keys():
            ws.cell(row=3, column=col).value = interval_name
            ws.cell(row=3, column=col).font = Font(bold=True, size=11)
            ws.cell(row=3, column=col).fill = PatternFill(start_color='CCE5FF', end_color='CCE5FF', fill_type='solid')
            ws.cell(row=3, column=col).alignment = Alignment(horizontal='center')
            col += 1

        # Собираем данные
        row = 4
        data_found = False

        for theatre_name in self.list_theatres():
            theatre = self.get_theatre(theatre_name)

            occupancy = {interval: {"occupied": 0, "total": 0} for interval in time_intervals.keys()}

            for hall in theatre["halls"]:
                for session in hall["sessions"]:
                    try:
                        session_time = datetime.strptime(session["start_time"], "%Y-%m-%d %H:%M")
                        hour = session_time.hour

                        # Определяем интервал
                        for interval_name, (start_h, end_h) in time_intervals.items():
                            if start_h < end_h:
                                if start_h <= hour < end_h:
                                    interval_key = interval_name
                                    break
                            else:  # Ночной интервал
                                if hour >= start_h or hour < end_h:
                                    interval_key = interval_name
                                    break
                        else:
                            continue

                        # Подсчитываем занятость
                        for row_seats in session["seats"]:
                            for seat in row_seats:
                                occupancy[interval_key]["total"] += 1
                                if seat:
                                    occupancy[interval_key]["occupied"] += 1

                        data_found = True
                    except ValueError:
                        pass

            # Записываем данные
            ws.cell(row=row, column=1).value = theatre_name
            ws.cell(row=row, column=1).font = Font(bold=True)

            col = 2
            for interval_name in time_intervals.keys():
                if occupancy[interval_name]["total"] > 0:
                    percentage = (occupancy[interval_name]["occupied"] / occupancy[interval_name]["total"]) * 100
                    ws.cell(row=row, column=col).value = round(percentage, 1)
                    ws.cell(row=row, column=col).number_format = '0.0"%"'
                else:
                    ws.cell(row=row, column=col).value = 0
                    ws.cell(row=row, column=col).number_format = '0.0"%"'

                ws.cell(row=row, column=col).alignment = Alignment(horizontal='center')
                col += 1

            row += 1

        if not data_found:
            ws['A5'] = 'Данные о сеансах не найдены'
            ws['A5'].font = Font(italic=True, color='999999')
        else:
            # Создаем диаграмму
            chart = BarChart()
            chart.type = "col"
            chart.style = 10
            chart.title = "Загруженность по времени суток (%)"
            chart.y_axis.title = 'Процент занятости'
            chart.x_axis.title = 'Время суток'

            data = Reference(ws, min_col=2, min_row=3, max_row=row - 1, max_col=5)
            cats = Reference(ws, min_col=2, min_row=3, max_row=3, max_col=5)

            chart.add_data(data, titles_from_data=True)
            chart.set_categories(cats)
            chart.height = 10
            chart.width = 20

            ws.add_chart(chart, f"A{row + 2}")

        # Автоподбор ширины столбцов
        for column in ws.columns:
            max_length = 0
            if isinstance(column[0], Cell):
                column_letter = column[0].column_letter
            else:
                column_letter = get_column_letter(column[0].column)
            for cell in column:
                try:
                    if len(str(cell.value)) > max_length:
                        max_length = len(str(cell.value))
                except:
                    pass
            adjusted_width = min(max_length + 2, 50)
            ws.column_dimensions[column_letter].width = adjusted_width

        # Сохраняем файл
        filename = os.path.join(self.reports_dir, f'occupancy_{datetime.now().strftime("%Y%m%d")}.xlsx')
        wb.save(filename)
        print(f"\n График загруженности успешно сохранен: {filename}")
        return filename

    def generate_movie_promo_pptx(self, movie_name):
        """Генерация рекламного буклета для фильма в формате PPTX"""
        print(f"\n--- ГЕНЕРАЦИЯ РЕКЛАМНОГО БУКЛЕТА ДЛЯ '{movie_name}' ---")

        prs = Presentation()
        prs.slide_width = PptxInches(10)
        prs.slide_height = PptxInches(7.5)

        # Слайд 1: Титульный
        slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Пустой слайд

        # Фон
        background = slide1.shapes.add_shape(
            1, 0, 0, prs.slide_width, prs.slide_height
        )
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = PptxRGBColor(10, 25, 47)

        # Заголовок
        title_box = slide1.shapes.add_textbox(
            PptxInches(1), PptxInches(2.5), PptxInches(8), PptxInches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = movie_name.upper()
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.alignment = PP_ALIGN.CENTER
        title_paragraph.font.size = PptxPt(60)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = PptxRGBColor(255, 215, 0)

        # Подзаголовок
        subtitle_box = slide1.shapes.add_textbox(
            PptxInches(1), PptxInches(4.5), PptxInches(8), PptxInches(0.8)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Смотрите в наших кинотеатрах!"
        subtitle_paragraph = subtitle_frame.paragraphs[0]
        subtitle_paragraph.alignment = PP_ALIGN.CENTER
        subtitle_paragraph.font.size = PptxPt(28)
        subtitle_paragraph.font.color.rgb = PptxRGBColor(255, 255, 255)

        # Собираем информацию о сеансах
        sessions_data = []

        for theatre_name in self.list_theatres():
            theatre = self.get_theatre(theatre_name)
            theatre_sessions = []

            for hall in theatre["halls"]:
                for session in hall["sessions"]:
                    if session["movie"] == movie_name:
                        try:
                            session_time = datetime.strptime(session["start_time"], "%Y-%m-%d %H:%M")
                            theatre_sessions.append({
                                "hall": hall["number"],
                                "time": session_time,
                                "duration": session["duration"]
                            })
                        except ValueError:
                            pass

            if theatre_sessions:
                theatre_sessions.sort(key=lambda x: x["time"])
                sessions_data.append({
                    "theatre": theatre_name,
                    "sessions": theatre_sessions
                })

        if not sessions_data:
            # Слайд с сообщением об отсутствии сеансов
            slide2 = prs.slides.add_slide(prs.slide_layouts[6])
            background2 = slide2.shapes.add_shape(
                1, 0, 0, prs.slide_width, prs.slide_height
            )
            fill2 = background2.fill
            fill2.solid()
            fill2.fore_color.rgb = PptxRGBColor(240, 240, 240)

            text_box = slide2.shapes.add_textbox(
                PptxInches(2), PptxInches(3), PptxInches(6), PptxInches(1)
            )
            text_frame = text_box.text_frame
            text_frame.text = f"Сеансы фильма '{movie_name}' не найдены"
            text_paragraph = text_frame.paragraphs[0]
            text_paragraph.alignment = PP_ALIGN.CENTER
            text_paragraph.font.size = PptxPt(24)
            text_paragraph.font.color.rgb = PptxRGBColor(100, 100, 100)
        else:
            # Слайды с информацией о кинотеатрах
            for theatre_data in sessions_data:
                slide = prs.slides.add_slide(prs.slide_layouts[6])

                # Фон
                background = slide.shapes.add_shape(
                    1, 0, 0, prs.slide_width, prs.slide_height
                )
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = PptxRGBColor(245, 245, 250)

                # Заголовок кинотеатра
                theatre_title_box = slide.shapes.add_textbox(
                    PptxInches(0.5), PptxInches(0.5), PptxInches(9), PptxInches(0.8)
                )
                theatre_title_frame = theatre_title_box.text_frame
                theatre_title_frame.text = f"🎬 {theatre_data['theatre']}"
                theatre_title_paragraph = theatre_title_frame.paragraphs[0]
                theatre_title_paragraph.alignment = PP_ALIGN.CENTER
                theatre_title_paragraph.font.size = PptxPt(36)
                theatre_title_paragraph.font.bold = True
                theatre_title_paragraph.font.color.rgb = PptxRGBColor(0, 51, 102)

                # Информация о сеансах
                sessions_box = slide.shapes.add_textbox(
                    PptxInches(1), PptxInches(1.8), PptxInches(8), PptxInches(5)
                )
                sessions_frame = sessions_box.text_frame
                sessions_frame.word_wrap = True

                for i, session in enumerate(theatre_data['sessions'][:10]):  # Максимум 10 сеансов на слайд
                    p = sessions_frame.add_paragraph() if i > 0 else sessions_frame.paragraphs[0]
                    p.text = f"📅 {session['time'].strftime('%d.%m.%Y')} в {session['time'].strftime('%H:%M')} | Зал {session['hall']} | {session['duration']} мин"
                    p.font.size = PptxPt(20)
                    p.font.color.rgb = PptxRGBColor(50, 50, 50)
                    p.space_after = PptxPt(12)

        # Сохраняем презентацию
        filename = os.path.join(self.reports_dir, f'promo_{movie_name.replace(" ", "_")}.pptx')
        prs.save(filename)
        print(f"\n✅ Рекламный буклет успешно сохранен: {filename}")
        return filename


def main():
    system = CinemaSystem()

    print("=" * 60)
    print("БИЛЕТНАЯ СИСТЕМА КИНОТЕАТРОВ".center(60))
    print("=" * 60)

    while True:
        print("\n" + "=" * 60)
        print("ГЛАВНОЕ МЕНЮ")
        print("=" * 60)
        print("1. Добавить кинотеатр")
        print("2. Добавить зал в кинотеатр")
        print("3. Создать сеанс")
        print("4. Продать билет")
        print("5. Найти ближайший сеанс фильма")
        print("6. Показать план зала")
        print("7. Показать список кинотеатров")
        print("8. Сформировать расписание сеансов за месяц (DOCX)")
        print("9. Сформировать график загруженности (XLSX)")
        print("10. Сформировать рекламный буклет фильма (PPTX)")
        print("0. Выход")
        print("=" * 60)

        choice = input("Выберите действие: ").strip()

        if choice == "1":
            print("\n--- ДОБАВЛЕНИЕ КИНОТЕАТРА ---")
            name = input("Введите название кинотеатра: ").strip()
            if name:
                system.add_theatre(name)
            else:
                print("Название не может быть пустым!")

        elif choice == "2":
            print("\n--- ДОБАВЛЕНИЕ ЗАЛА ---")
            theatres = system.list_theatres()
            if not theatres:
                print("Сначала добавьте хотя бы один кинотеатр!")
                continue

            print("Доступные кинотеатры:", ", ".join(theatres))
            theatre_name = input("Введите название кинотеатра: ").strip()

            try:
                hall_number = int(input("Введите номер зала: "))
                rows = int(input("Введите количество рядов: "))
                seats = int(input("Введите количество мест в ряду: "))
                system.add_hall(theatre_name, hall_number, rows, seats)
            except ValueError:
                print("Ошибка! Введите числовые значения.")

        elif choice == "3":
            print("\n--- СОЗДАНИЕ СЕАНСА ---")
            theatres = system.list_theatres()
            if not theatres:
                print("Сначала добавьте хотя бы один кинотеатр!")
                continue

            print("Доступные кинотеатры:", ", ".join(theatres))
            theatre_name = input("Введите название кинотеатра: ").strip()

            theatre = system.get_theatre(theatre_name)
            if not theatre:
                print(f"Кинотеатр '{theatre_name}' не найден!")
                continue

            if not theatre["halls"]:
                print("В этом кинотеатре нет залов!")
                continue

            print("Доступные залы:", ", ".join([str(h["number"]) for h in theatre["halls"]]))

            try:
                hall_number = int(input("Введите номер зала: "))
                movie = input("Введите название фильма: ").strip()
                start_time = input("Введите время начала (ГГГГ-ММ-ДД ЧЧ:ММ): ").strip()
                duration = int(input("Введите длительность фильма (в минутах): "))
                system.create_session(theatre_name, hall_number, movie, start_time, duration)
            except ValueError:
                print("Ошибка! Проверьте формат введённых данных.")

        elif choice == "4":
            print("\n--- ПРОДАЖА БИЛЕТА ---")
            theatres = system.list_theatres()
            if not theatres:
                print("Сначала добавьте хотя бы один кинотеатр!")
                continue

            print("Доступные кинотеатры:", ", ".join(theatres))
            theatre_name = input("Введите название кинотеатра: ").strip()

            theatre = system.get_theatre(theatre_name)
            if not theatre:
                print(f"Кинотеатр '{theatre_name}' не найден!")
                continue

            if not theatre["halls"]:
                print("В этом кинотеатре нет залов!")
                continue

            print("Доступные залы:", ", ".join([str(h["number"]) for h in theatre["halls"]]))

            try:
                hall_number = int(input("Введите номер зала: "))

                hall = None
                for h in theatre["halls"]:
                    if h["number"] == hall_number:
                        hall = h
                        break

                if not hall:
                    print(f"Зал №{hall_number} не найден!")
                    continue

                if not hall["sessions"]:
                    print("В этом зале нет сеансов!")
                    continue

                print("\nСеансы:")
                for i, session in enumerate(hall["sessions"]):
                    print(f"{i}. {session['movie']} - {session['start_time']}")

                session_index = int(input("Введите номер сеанса: "))
                system.print_hall_plan(theatre_name, hall_number, session_index)

                row = int(input("Введите номер ряда: ")) - 1
                seat = int(input("Введите номер места: ")) - 1

                system.sell_ticket(theatre_name, hall_number, session_index, row, seat)
            except ValueError:
                print("Ошибка! Введите числовые значения.")

        elif choice == "5":
            print("\n--- ПОИСК БЛИЖАЙШЕГО СЕАНСА ---")
            movie = input("Введите название фильма: ").strip()
            if movie:
                system.find_nearest_session(movie)
            else:
                print("Название фильма не может быть пустым!")

        elif choice == "6":
            print("\n--- ПЛАН ЗАЛА ---")
            theatres = system.list_theatres()
            if not theatres:
                print("Сначала добавьте хотя бы один кинотеатр!")
                continue

            print("Доступные кинотеатры:", ", ".join(theatres))
            theatre_name = input("Введите название кинотеатра: ").strip()

            theatre = system.get_theatre(theatre_name)
            if not theatre:
                print(f"Кинотеатр '{theatre_name}' не найден!")
                continue

            if not theatre["halls"]:
                print("В этом кинотеатре нет залов!")
                continue

            print("Доступные залы:", ", ".join([str(h["number"]) for h in theatre["halls"]]))

            try:
                hall_number = int(input("Введите номер зала: "))

                hall = None
                for h in theatre["halls"]:
                    if h["number"] == hall_number:
                        hall = h
                        break

                if not hall:
                    print(f"Зал №{hall_number} не найден!")
                    continue

                if not hall["sessions"]:
                    print("В этом зале нет сеансов!")
                    continue

                print("\nСеансы:")
                for i, session in enumerate(hall["sessions"]):
                    print(f"{i}. {session['movie']} - {session['start_time']}")

                session_index = int(input("Введите номер сеанса: "))
                system.print_hall_plan(theatre_name, hall_number, session_index)
            except ValueError:
                print("Ошибка! Введите числовые значения.")

        elif choice == "7":
            print("\n--- СПИСОК КИНОТЕАТРОВ ---")
            theatres = system.list_theatres()
            if not theatres:
                print("Кинотеатры не найдены.")
            else:
                for theatre_name in theatres:
                    theatre = system.get_theatre(theatre_name)
                    print(f"\nКинотеатр: {theatre_name}")
                    print(f"Количество залов: {len(theatre['halls'])}")
                    for hall in theatre["halls"]:
                        print(f"  Зал №{hall['number']}: {hall['rows']} рядов x {hall['seats_per_row']} мест, "
                              f"Сеансов: {len(hall['sessions'])}")

        elif choice == "8":
            try:
                system.generate_monthly_schedule_docx()
            except Exception as e:
                print(f"Ошибка при генерации расписания: {e}")

        elif choice == "9":
            try:
                system.generate_occupancy_chart_xlsx()
            except Exception as e:
                print(f"Ошибка при генерации графика: {e}")

        elif choice == "10":
            movie_name = input("Введите название фильма для буклета: ").strip()
            if movie_name:
                try:
                    system.generate_movie_promo_pptx(movie_name)
                except Exception as e:
                    print(f"Ошибка при генерации буклета: {e}")
            else:
                print("Название фильма не может быть пустым!")

        elif choice == "0":
            print("\nСпасибо за использование билетной системы! До свидания!")
            break

        else:
            print("Неверный выбор! Попробуйте снова.")


if __name__ == "__main__":
    main()
